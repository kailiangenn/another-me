"""
AME统一解析器

整合所有解析方法，用户继承时只需关注特定格式的解析实现
"""

from pathlib import Path
from typing import Dict, Any
from loguru import logger
import PyPDF2

from .base import FileParserBase
from ..utils.models import DocumentFormat


class AMEParser(FileParserBase):
    """
    AME统一文档解析器
    
    已实现的格式：
    - PDF (PyPDF2)
    - DOCX (python-docx)
    - PPT/PPTX (python-pptx)
    - Markdown
    - Text
    
    用户继承使用：
    1. 只需继承 AMEParser
    2. 重写感兴趣的 _parse_* 方法
    3. 无需关注主流程和Markdown转换
    
    示例：
        class MyParser(AMEParser):
            async def _parse_pdf(self, path: Path):
                # 自定义PDF解析逻辑
                content = "..."
                metadata = {...}
                return content, metadata
    """
    
    SUPPORTED_EXTENSIONS = {
        "pdf", "docx", "doc", "ppt", "pptx",
        "md", "markdown", "txt", "text"
    }
    
    def __init__(self, use_pdfplumber: bool = False):
        """
        初始化AME解析器
        
        Args:
            use_pdfplumber: PDF解析是否使用pdfplumber（默认PyPDF2）
        """
        super().__init__()
        self.use_pdfplumber = use_pdfplumber
    
    def can_parse(self, file_path: str) -> bool:
        """判断是否支持该文件"""
        extension = self._get_file_extension(file_path)
        return extension in self.SUPPORTED_EXTENSIONS
        
    async def _parse_pdf(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """
        解析PDF文件
        
        Args:
            path: 文件路径
        
        Returns:
            (content, metadata): 内容和元数据
        """
        if self.use_pdfplumber:
            return await self._parse_pdf_with_pdfplumber(path)
        else:
            return await self._parse_pdf_with_pypdf2(path)
    
    async def _parse_pdf_with_pypdf2(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """使用PyPDF2解析PDF"""
        content_parts = []
        
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            total_pages = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    content_parts.append(text)
        
        content = "\n\n".join(content_parts)
        
        metadata = {
            "parser": "PyPDF2",
            "file_size": path.stat().st_size,
            "total_pages": total_pages,
        }
        
        return content, metadata
    
    async def _parse_pdf_with_pdfplumber(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """使用pdfplumber解析PDF（质量更高）"""
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber未安装，回退到PyPDF2")
            return await self._parse_pdf_with_pypdf2(path)
        
        content_parts = []
        
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    content_parts.append(text)
        
        content = "\n\n".join(content_parts)
        
        metadata = {
            "parser": "pdfplumber",
            "file_size": path.stat().st_size,
            "total_pages": total_pages,
        }
        
        return content, metadata
        
    async def _parse_docx(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """
        解析DOCX文件
        
        Args:
            path: 文件路径
        
        Returns:
            (content, metadata): 内容和元数据
        """
        try:
            import docx
        except ImportError:
            raise ImportError("请安装python-docx: pip install python-docx")
        
        doc = docx.Document(path)
        content_parts = []
        sections = []  # 保存结构化信息
        
        # 提取段落（保留样式信息）
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            content_parts.append(text)
            
            # 判断是否是标题
            if para.style.name.startswith("Heading"):
                # 提取标题级别
                level = 1
                try:
                    level_str = para.style.name.replace("Heading", "").strip()
                    if level_str.isdigit():
                        level = int(level_str)
                except:
                    pass
                
                sections.append({
                    "type": f"h{level}",
                    "content": text,
                    "style": para.style.name
                })
            else:
                sections.append({
                    "type": "paragraph",
                    "content": text,
                    "style": para.style.name
                })
        
        # 提取表格
        for table in doc.tables:
            table_text = self._extract_docx_table(table)
            if table_text:
                content_parts.append(table_text)
                sections.append({
                    "type": "table",
                    "content": table_text
                })
        
        content = "\n\n".join(content_parts)
        
        # 提取文档属性
        core_props = doc.core_properties
        metadata = {
            "parser": "python-docx",
            "file_size": path.stat().st_size,
            "author": core_props.author if core_props.author else None,
            "title": core_props.title if core_props.title else None,
            "sections": sections,  # 添加结构化信息
        }
        
        return content, metadata
    
    def _extract_docx_table(self, table) -> str:
        """提取DOCX表格"""
        lines = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(" | ".join(cells))
        return "\n".join(lines)
    
    async def _parse_doc(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """
        解析DOC文件（旧版Word格式）
        
        注意：DOC格式解析较复杂，建议转换为DOCX
        """
        logger.warning(f"DOC格式支持有限，建议转换为DOCX: {path.name}")
        
        # 尝试使用docx库（可能支持部分DOC文件）
        try:
            import docx
            doc = docx.Document(path)
            content_parts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
            content = "\n\n".join(content_parts)
            
            metadata = {
                "parser": "python-docx",
                "file_size": path.stat().st_size,
                "warning": "DOC格式支持有限"
            }
            
            return content, metadata
        
        except Exception as e:
            logger.error(f"DOC解析失败: {e}")
            raise NotImplementedError(f"DOC格式解析失败，请转换为DOCX格式: {path.name}")
    
    async def _parse_ppt(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """
        解析PPT/PPTX文件
        
        Args:
            path: 文件路径
        
        Returns:
            (content, metadata): 内容和元数据
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装python-pptx: pip install python-pptx")
        
        prs = Presentation(path)
        content_parts = []
        sections = []  # 保存结构化信息
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            slide_sections = []
            
            for shape in slide.shapes:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue
                
                text = shape.text.strip()
                slide_texts.append(text)
                
                # 判断是否是标题占位符
                section_type = "paragraph"
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    placeholder = shape.placeholder_format
                    if placeholder.type == 1:  # Title placeholder
                        section_type = "h2"
                
                slide_sections.append({
                    "type": section_type,
                    "content": text,
                    "slide_num": slide_num
                })
            
            if slide_texts:
                content_parts.append(
                    f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts)
                )
                sections.extend(slide_sections)
        
        content = "\n\n".join(content_parts)
        
        metadata = {
            "parser": "python-pptx",
            "file_size": path.stat().st_size,
            "slide_count": len(prs.slides),
            "sections": sections,  # 添加结构化信息
        }
        
        return content, metadata
    
    async def _parse_markdown(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """
        解析Markdown文件
        
        Args:
            path: 文件路径
        
        Returns:
            (content, metadata): 内容和元数据
        """
        content = await self._read_text_file(path)
        
        metadata = {
            "parser": "markdown",
            "file_size": path.stat().st_size,
        }
        
        return content, metadata
    
    # ==================== Text解析 ====================
    
    async def _parse_text(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """
        解析纯文本文件
        
        Args:
            path: 文件路径
        
        Returns:
            (content, metadata): 内容和元数据
        """
        content = await self._read_text_file(path)
        
        metadata = {
            "parser": "text",
            "file_size": path.stat().st_size,
        }
        
        return content, metadata
    
    async def _read_text_file(self, path: Path) -> str:
        """
        读取文本文件（自动处理编码）
        
        Args:
            path: 文件路径
        
        Returns:
            content: 文件内容
        """
        encodings = ["utf-8", "gbk", "gb2312", "iso-8859-1"]
        
        for encoding in encodings:
            try:
                with open(path, "r", encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, LookupError):
                continue
        
        # 如果所有编码都失败，使用utf-8并忽略错误
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            logger.warning(f"使用 utf-8 with errors='ignore' 读取文件: {path}")
            return f.read()

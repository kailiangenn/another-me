"""
PPT/PPTX 文件解析器

支持:
- .ppt/.pptx 文件
- 幻灯片文本提取
- 标题识别
- 列表提取

依赖: python-pptx
"""

from typing import List
from pathlib import Path
from loguru import logger
from pptx import Presentation

from .base import FileParserBase
from ..core.models import ParsedDocument, DocumentSection, DocumentFormat, SectionType
from ..core.exceptions import ParseError


class PPTParser(FileParserBase):
    """
    PPT/PPTX 解析器
    
    特性:
    - 提取每页幻灯片内容
    - 识别标题和正文
    - 提取列表项
    - 保留页码信息
    
    依赖:
    - python-pptx
    """
    
    SUPPORTED_EXTENSIONS = {"ppt", "pptx"}
    
    def can_parse(self, file_path: str) -> bool:
        """判断是否支持该文件"""
        extension = self._get_file_extension(file_path)
        return extension in self.SUPPORTED_EXTENSIONS
    
    async def parse(self, file_path: str) -> ParsedDocument:
        """
        解析 PPT/PPTX 文件
        
        Args:
            file_path: 文件路径
        
        Returns:
            parsed_doc: 解析结果
        """
        path = self._validate_file_exists(file_path)
        
        try:
            return await self._parse_pptx(path)
        except Exception as e:
            logger.error(f"PPT 文件解析失败: {file_path}, 错误: {e}")
            raise
    
    async def _parse_pptx(self, path: Path) -> ParsedDocument:
        """
        解析 PPTX 文档
        
        Args:
            path: 文件路径
        
        Returns:
            parsed_doc: 解析结果
        """
        prs = Presentation(path)
        
        sections = []
        raw_content_parts = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            
            for shape in slide.shapes:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue
                
                text = shape.text.strip()
                slide_texts.append(text)
                
                # 判断是否是标题
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    placeholder = shape.placeholder_format
                    if placeholder.type == 1:  # Title placeholder
                        sections.append(DocumentSection(
                            type=SectionType.HEADING_1,
                            content=text,
                            level=1,
                            page_number=slide_num,
                            metadata={"slide_num": slide_num, "shape_type": "title"}
                        ))
                        continue
                
                # 检查是否包含列表
                if self._is_list_text(text):
                    for line in text.split('\n'):
                        line = line.strip()
                        if line and (line.startswith(('•', '-', '*', '○', '■')) or 
                                   (len(line) > 2 and line[0].isdigit() and line[1] in '.)')):
                            sections.append(DocumentSection(
                                type=SectionType.LIST_ITEM,
                                content=line.lstrip('•-*○■ ').lstrip('0123456789.)').strip(),
                                page_number=slide_num,
                                metadata={"slide_num": slide_num}
                            ))
                else:
                    # 普通段落
                    sections.append(DocumentSection(
                        type=SectionType.PARAGRAPH,
                        content=text,
                        page_number=slide_num,
                        metadata={"slide_num": slide_num}
                    ))
            
            if slide_texts:
                raw_content_parts.append(f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        
        raw_content = "\n\n".join(raw_content_parts)
        
        return ParsedDocument(
            format=DocumentFormat.PPT,
            file_path=str(path),
            raw_content=raw_content,
            sections=sections,
            total_pages=len(prs.slides),
            metadata={
                "parser": "python-pptx",
                "file_size": path.stat().st_size,
                "slide_count": len(prs.slides),
            }
        )
    
    def _is_list_text(self, text: str) -> bool:
        """
        判断文本是否包含列表
        
        Args:
            text: 文本内容
        
        Returns:
            is_list: 是否是列表
        """
        lines = text.split('\n')
        if len(lines) < 2:
            return False
        
        list_marker_count = 0
        for line in lines:
            line = line.strip()
            if line.startswith(('•', '-', '*', '○', '■')):
                list_marker_count += 1
            elif len(line) > 2 and line[0].isdigit() and line[1] in '.)':
                list_marker_count += 1
        
        return list_marker_count >= 2
